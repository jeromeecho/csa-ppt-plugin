#!/usr/bin/env python3
"""Auto-fix text overflow issues in PPTX files.

Priority: 内容完整显示 > 不换行 > 字体大小
(Content fully visible > no wrapping > font size)

Fix strategies applied in order:
1. Expand container (if space available, avoids wrapping)
2. Widen shape (eliminate widow lines without shrinking font)
3. Reduce font size (down to 9pt minimum, prevents wrapping)
4. Reduce paragraph spacing
5. Truncate text as last resort

Usage:
    python fix_overflow.py input.pptx --report validation.json -o output.pptx
    python fix_overflow.py input.pptx -o output.pptx  # auto-detect issues
"""

import argparse
import json
import os
import sys
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("ERROR: Pillow is required. Install with: pip install Pillow", file=sys.stderr)
    sys.exit(2)

# Import validation functions
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from validate_visual import (
    measure_text_height, load_font, wrap_text_line, emu_to_inches,
    validate_pptx, detect_widow_lines, ValidationIssue
)


# ─── Fix Strategies ──────────────────────────────────────────────────────────

MIN_FONT_SIZE_PT = 7  # Content visibility is the hard constraint; font size is soft
MIN_SPACING_PT = 0


def get_shape_by_index(slide, shape_idx: int):
    """Get shape from slide by index."""
    shapes = list(slide.shapes)
    if 0 <= shape_idx < len(shapes):
        return shapes[shape_idx]
    return None


def get_current_font_sizes(text_frame) -> List[float]:
    """Get all font sizes used in a text frame."""
    sizes = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
    return sizes


def reduce_font_size(text_frame, reduction_pt: float = 1.0) -> bool:
    """Reduce all font sizes in the text frame by reduction_pt.

    Returns True if any reduction was applied (above minimum).
    """
    changed = False
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                current = run.font.size.pt
                new_size = max(MIN_FONT_SIZE_PT, current - reduction_pt)
                if new_size < current:
                    run.font.size = Pt(new_size)
                    changed = True
            else:
                # No explicit size - set to a reasonable default minus reduction
                run.font.size = Pt(max(MIN_FONT_SIZE_PT, 12 - reduction_pt))
                changed = True
    return changed


def reduce_spacing(text_frame) -> bool:
    """Reduce paragraph spacing (space_before, space_after)."""
    changed = False
    for para in text_frame.paragraphs:
        if para.space_before and para.space_before > Pt(MIN_SPACING_PT):
            new_val = max(Pt(MIN_SPACING_PT), para.space_before - Pt(1))
            if new_val != para.space_before:
                para.space_before = new_val
                changed = True
        if para.space_after and para.space_after > Pt(MIN_SPACING_PT):
            new_val = max(Pt(MIN_SPACING_PT), para.space_after - Pt(1))
            if new_val != para.space_after:
                para.space_after = new_val
                changed = True
    return changed


def can_expand_shape(shape, slide, slide_width_emu, slide_height_emu,
                     expand_inches: float = 0.1) -> bool:
    """Check if shape can be expanded downward without exceeding slide bounds."""
    new_bottom = shape.top + shape.height + Inches(expand_inches)
    if new_bottom > slide_height_emu:
        return False

    # Check overlap with other shapes below
    expand_bottom = shape.top + shape.height + Inches(expand_inches)
    for other in slide.shapes:
        if other == shape:
            continue
        # Check if other shape is below and would overlap
        if other.top >= shape.top + shape.height:
            if other.top < expand_bottom:
                # Check horizontal overlap
                shape_right = shape.left + shape.width
                other_right = other.left + other.width
                if not (shape_right <= other.left or shape.left >= other_right):
                    return False
    return True


def expand_shape_height(shape, expand_inches: float = 0.1) -> None:
    """Expand shape height by the given amount."""
    shape.height = shape.height + Inches(expand_inches)


def truncate_text(text_frame, shape_width_inches: float, shape_height_inches: float) -> bool:
    """Truncate text to fit within the shape. Last resort."""
    changed = False
    paragraphs = text_frame.paragraphs

    # Try removing paragraphs from the end until it fits
    # We can't actually remove paragraphs easily in python-pptx,
    # but we can clear runs in the last paragraphs
    while len(paragraphs) > 1:
        total_h, overflow = measure_text_height(text_frame, shape_width_inches, shape_height_inches)
        if overflow <= 0.05:
            break

        # Find the last paragraph with text
        last_para = None
        for para in reversed(paragraphs):
            if para.text.strip():
                last_para = para
                break

        if last_para is None:
            break

        # Truncate the last paragraph
        if last_para.runs:
            last_run = last_para.runs[-1]
            text = last_run.text
            if len(text) > 3:
                last_run.text = text[:len(text)//2] + "..."
                changed = True
            else:
                # Clear this run entirely
                last_run.text = ""
                changed = True
        else:
            break

    return changed


def fix_widow_line(shape, slide, slide_width_emu) -> Optional[str]:
    """Fix widow lines by reducing font or expanding width.

    A widow line is when text wraps with the last line containing < 20% of content.
    Fix: reduce font by 0.5pt increments until text fits on fewer lines,
    or widen the shape if possible.
    """
    if not shape.has_text_frame:
        return None

    tf = shape.text_frame
    shape_width = emu_to_inches(shape.width)
    shape_height = emu_to_inches(shape.height)

    # Check if there are widow issues
    widows = detect_widow_lines(tf, shape_width, shape_height)
    if not widows:
        return None

    slide_width_inches = emu_to_inches(slide_width_emu)

    # Strategy 1: Try widening the shape (up to slide edge)
    shape_left = emu_to_inches(shape.left) if shape.left else 0
    available_right = slide_width_inches - shape_left - 0.1  # 0.1" margin from edge
    max_expand = available_right - shape_width

    if max_expand > 0.05:
        # Try expanding width by small increments
        from pptx.util import Inches as PtxInches
        for expand in [0.1, 0.2, 0.3, 0.4, 0.5]:
            if expand > max_expand:
                break
            test_width = shape_width + expand
            test_widows = detect_widow_lines(tf, test_width, shape_height)
            if not test_widows:
                shape.width = shape.width + PtxInches(expand)
                return f"widened_shape (+{expand:.2f}\")"

    # Strategy 2: Reduce font size by 0.5pt increments (max 3 rounds)
    for attempt in range(3):
        if reduce_font_size(tf, reduction_pt=0.5):
            widows = detect_widow_lines(tf, shape_width, shape_height)
            if not widows:
                sizes = get_current_font_sizes(tf)
                min_size = min(sizes) if sizes else 0
                return f"reduced_font_for_widow (min now {min_size:.1f}pt)"
        else:
            break

    return None


# ─── Main Fix Engine ─────────────────────────────────────────────────────────

def fix_text_overflow(shape, slide, slide_width_emu, slide_height_emu) -> Optional[str]:
    """Fix text overflow for a single shape. Returns description of fix applied."""
    if not shape.has_text_frame:
        return None

    tf = shape.text_frame
    shape_width = emu_to_inches(shape.width)
    shape_height = emu_to_inches(shape.height)

    # Check if there's actually overflow
    total_h, overflow = measure_text_height(tf, shape_width, shape_height)
    if overflow <= 0.05:
        return None

    # Priority: 内容完整 > 不换行 > 字体大小
    # Strategy 1: Expand container FIRST (preserves font size, avoids wrapping)
    needed = overflow + 0.05
    if can_expand_shape(shape, slide, slide_width_emu, slide_height_emu, expand_inches=needed):
        expand_shape_height(shape, needed)
        return f"expanded_height (+{needed:.2f}\")"

    # Strategy 2: Reduce font size (up to 6 rounds of 0.5pt each, min 9pt)
    for _ in range(6):
        if reduce_font_size(tf, reduction_pt=0.5):
            _, overflow = measure_text_height(tf, shape_width, shape_height)
            if overflow <= 0.05:
                sizes = get_current_font_sizes(tf)
                min_size = min(sizes) if sizes else 0
                return f"reduced_font (min now {min_size:.1f}pt)"
        else:
            break  # Can't reduce further (at minimum)

    # Strategy 3: Reduce spacing
    for _ in range(3):
        if reduce_spacing(tf):
            _, overflow = measure_text_height(tf, shape_width, shape_height)
            if overflow <= 0.05:
                return "reduced_spacing"
        else:
            break

    # Strategy 4: Last attempt — expand even for small remaining overflow
    _, overflow = measure_text_height(tf, shape_width, shape_height)
    if overflow > 0:
        needed = overflow + 0.05
        if can_expand_shape(shape, slide, slide_width_emu, slide_height_emu, expand_inches=needed):
            expand_shape_height(shape, needed)
            return f"expanded_height (+{needed:.2f}\")"

    # Strategy 4: Truncate as last resort
    if truncate_text(tf, shape_width, shape_height):
        return "truncated_text"

    return None


def fix_font_too_small(shape) -> Optional[str]:
    """Fix font sizes that are below minimum."""
    if not shape.has_text_frame:
        return None

    changed = False
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size and run.font.size.pt < MIN_FONT_SIZE_PT:
                run.font.size = Pt(MIN_FONT_SIZE_PT)
                changed = True

    return "increased_to_10pt" if changed else None


# ─── Main ────────────────────────────────────────────────────────────────────

def fix_pptx(input_path: str, output_path: str,
             report_path: Optional[str] = None, threshold: float = 0.05) -> dict:
    """Fix all overflow issues in a PPTX file.

    Args:
        input_path: Source PPTX
        output_path: Where to save fixed PPTX
        report_path: Optional validation report JSON (if None, auto-detects issues)
        threshold: Overflow threshold in inches

    Returns:
        Summary dict with fix counts and details
    """
    prs = Presentation(input_path)
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    # Get issues from report or auto-detect
    if report_path and os.path.isfile(report_path):
        with open(report_path, "r", encoding="utf-8") as f:
            report = json.load(f)
        issues = report.get("issues", [])
    else:
        # Auto-detect
        raw_issues = validate_pptx(input_path, threshold=threshold)
        issues = [i.to_dict() for i in raw_issues]

    if not issues:
        # No issues, just copy
        prs.save(output_path)
        return {"status": "no_issues", "fixes_applied": 0, "details": []}

    fixes_applied = []
    slides_fixed = set()

    # Group issues by slide
    by_slide = {}
    for issue in issues:
        slide_num = issue["slide"]
        by_slide.setdefault(slide_num, []).append(issue)

    for slide_num, slide_issues in sorted(by_slide.items()):
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]

        # Process each issue
        processed_shapes = set()
        for issue in slide_issues:
            shape_idx = issue["shape_index"]
            issue_type = issue["type"]

            # Don't fix same shape twice
            if (slide_num, shape_idx) in processed_shapes:
                continue

            shape = get_shape_by_index(slide, shape_idx)
            if shape is None:
                continue

            fix_desc = None
            if issue_type == "text_overflow":
                fix_desc = fix_text_overflow(shape, slide, slide_width_emu, slide_height_emu)
            elif issue_type == "font_too_small":
                fix_desc = fix_font_too_small(shape)
            elif issue_type == "widow_line":
                fix_desc = fix_widow_line(shape, slide, slide_width_emu)

            if fix_desc:
                processed_shapes.add((slide_num, shape_idx))
                slides_fixed.add(slide_num)
                fixes_applied.append({
                    "slide": slide_num,
                    "shape_index": shape_idx,
                    "issue_type": issue_type,
                    "fix": fix_desc,
                    "text_preview": issue.get("text_preview", "")[:40],
                })

    prs.save(output_path)

    return {
        "status": "fixed",
        "fixes_applied": len(fixes_applied),
        "slides_affected": len(slides_fixed),
        "details": fixes_applied,
    }


def main():
    parser = argparse.ArgumentParser(description="Auto-fix text overflow in PPTX files")
    parser.add_argument("input", help="Path to .pptx file to fix")
    parser.add_argument("--report", "-r", help="Validation report JSON (from validate_visual.py)")
    parser.add_argument("--output", "-o", required=True, help="Output .pptx path")
    parser.add_argument("--threshold", type=float, default=0.05,
                        help="Overflow threshold in inches (default: 0.05)")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"ERROR: File not found: {args.input}", file=sys.stderr)
        sys.exit(2)

    print(f"Fixing: {args.input}")
    result = fix_pptx(args.input, args.output, args.report, args.threshold)

    print(f"RESULT: {result['status']}")
    print(f"FIXES_APPLIED: {result['fixes_applied']}")

    if result['details']:
        print(f"SLIDES_AFFECTED: {result.get('slides_affected', 0)}")
        print("\nFix details:")
        for fix in result['details']:
            print(f"  Slide {fix['slide']}, shape {fix['shape_index']}: "
                  f"{fix['issue_type']} → {fix['fix']}")
            if fix['text_preview']:
                print(f"    Text: \"{fix['text_preview']}...\"")

    print(f"\nOUTPUT_FILE: {args.output}")


if __name__ == "__main__":
    main()
