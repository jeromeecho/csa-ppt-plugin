#!/usr/bin/env python3
"""Validate PPTX for text overflow and layout issues.

Uses PIL-based text measurement to detect text that overflows its container.
Exit code 0 = all slides pass, exit code 1 = issues found.

Usage:
    python validate_visual.py input.pptx [--output report.json] [--threshold 0.05]
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("ERROR: Pillow is required. Install with: pip install Pillow", file=sys.stderr)
    sys.exit(2)


# ─── Font Resolution ─────────────────────────────────────────────────────────

FONT_SEARCH_PATHS = [
    # macOS
    "/System/Library/Fonts/",
    "/Library/Fonts/",
    os.path.expanduser("~/Library/Fonts/"),
    # Linux
    "/usr/share/fonts/",
    "/usr/local/share/fonts/",
    os.path.expanduser("~/.fonts/"),
    # Windows
    "C:/Windows/Fonts/",
]

FONT_NAME_MAP = {
    "arial": ["Arial.ttf", "arial.ttf", "Arial Unicode.ttf"],
    "microsoft yahei": ["msyh.ttc", "msyh.ttf", "MicrosoftYaHei.ttf", "Microsoft YaHei.ttf"],
    "segoe ui": ["segoeui.ttf", "SegoeUI.ttf"],
    "courier new": ["cour.ttf", "CourierNew.ttf"],
    "calibri": ["calibri.ttf", "Calibri.ttf"],
    "times new roman": ["times.ttf", "TimesNewRoman.ttf"],
}

_font_cache = {}


def find_font_path(font_name: str) -> Optional[str]:
    """Find the system path for a font by name."""
    if font_name in _font_cache:
        return _font_cache[font_name]

    key = font_name.lower().strip()
    candidates = FONT_NAME_MAP.get(key, [f"{font_name}.ttf", f"{font_name}.ttc"])

    for search_dir in FONT_SEARCH_PATHS:
        if not os.path.isdir(search_dir):
            continue
        for root, _, files in os.walk(search_dir):
            for candidate in candidates:
                if candidate in files:
                    path = os.path.join(root, candidate)
                    _font_cache[font_name] = path
                    return path

    _font_cache[font_name] = None
    return None


def load_font(font_name: str, font_size: int) -> ImageFont.FreeTypeFont:
    """Load a PIL font, falling back to default if not found."""
    path = find_font_path(font_name)
    if path:
        try:
            return ImageFont.truetype(path, size=font_size)
        except Exception:
            pass
    # Fallback: try common system font
    for fallback in ["Arial", "Helvetica", "DejaVuSans"]:
        path = find_font_path(fallback)
        if path:
            try:
                return ImageFont.truetype(path, size=font_size)
            except Exception:
                continue
    return ImageFont.load_default()


# ─── Text Measurement ─────────────────────────────────────────────────────────

def emu_to_inches(emu: int) -> float:
    """Convert EMU to inches."""
    return emu / 914400.0


def wrap_text_line(line: str, max_width_px: int, draw: ImageDraw.ImageDraw, font) -> List[str]:
    """Wrap a single line of text to fit within max_width_px."""
    if not line:
        return [""]

    try:
        line_width = draw.textlength(line, font=font)
    except Exception:
        line_width = len(line) * 8  # rough fallback

    if line_width <= max_width_px:
        return [line]

    # Word-wrap for space-separated text
    wrapped = []
    words = line.split(" ")
    if len(words) <= 1:
        # Single long word (like a URL or CJK without spaces) - character-level wrap
        chars = list(line)
        current = ""
        for ch in chars:
            test = current + ch
            try:
                w = draw.textlength(test, font=font)
            except Exception:
                w = len(test) * 8
            if w <= max_width_px:
                current = test
            else:
                if current:
                    wrapped.append(current)
                current = ch
        if current:
            wrapped.append(current)
        return wrapped if wrapped else [line]

    current_line = ""
    for word in words:
        test_line = current_line + (" " if current_line else "") + word
        try:
            w = draw.textlength(test_line, font=font)
        except Exception:
            w = len(test_line) * 8
        if w <= max_width_px:
            current_line = test_line
        else:
            if current_line:
                wrapped.append(current_line)
            current_line = word

    if current_line:
        wrapped.append(current_line)

    return wrapped if wrapped else [line]


def measure_text_height(text_frame, shape_width_inches: float, shape_height_inches: float) -> Tuple[float, float]:
    """Measure the total text height and compare against available height.

    Returns (total_height_inches, overflow_inches).
    overflow_inches > 0 means text overflows.

    IMPORTANT: PIL text measurement consistently UNDERESTIMATES compared to
    PowerPoint's actual rendering. PowerPoint uses wider line spacing, additional
    paragraph spacing, and different font metrics. We apply a scaled safety factor
    to compensate for this known discrepancy — larger for dense multi-line text
    (where accumulated error is significant) and minimal for short text.
    """
    # Safety factor scales with paragraph count:
    # - 1-2 paragraphs: 1.05x (titles, short labels — minimal accumulated error)
    # - 3-5 paragraphs: 1.15x (moderate content)
    # - 6+ paragraphs: 1.25x (dense content — accumulated spacing differences)
    # This prevents false positives on titles while catching real overflow on dense blocks.

    # Get margins
    margin_top = emu_to_inches(text_frame.margin_top) if text_frame.margin_top else 0.05
    margin_bottom = emu_to_inches(text_frame.margin_bottom) if text_frame.margin_bottom else 0.05
    margin_left = emu_to_inches(text_frame.margin_left) if text_frame.margin_left else 0.1
    margin_right = emu_to_inches(text_frame.margin_right) if text_frame.margin_right else 0.1

    usable_width = shape_width_inches - margin_left - margin_right
    usable_height = shape_height_inches - margin_top - margin_bottom

    if usable_width <= 0 or usable_height <= 0:
        return (0, 0)

    usable_width_px = int(usable_width * 96)
    usable_height_px = int(usable_height * 96)

    # Set up PIL measurement
    dummy_img = Image.new("RGB", (1, 1))
    draw = ImageDraw.Draw(dummy_img)

    total_height_px = 0.0

    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        text = paragraph.text
        if not text.strip():
            # Empty paragraph still takes some space (PowerPoint ~0.5 line height)
            total_height_px += 10  # more realistic for PowerPoint empty paragraphs
            continue

        # Get font info from first run
        font_name = "Arial"
        font_size = 14  # default pt
        if paragraph.runs:
            run = paragraph.runs[0]
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = int(run.font.size.pt)
        elif paragraph.font and paragraph.font.size:
            font_size = int(paragraph.font.size.pt)

        font = load_font(font_name, font_size)

        # Wrap text
        all_wrapped_lines = []
        for line in text.split("\n"):
            wrapped = wrap_text_line(line, usable_width_px, draw, font)
            all_wrapped_lines.extend(wrapped)

        if all_wrapped_lines:
            # Line height: font_size * standard single-line spacing (1.2x)
            # PowerPoint "single" spacing = font_size * 1.2
            line_height_px = font_size * 96 / 72 * 1.2

            # Check for custom line spacing
            if hasattr(paragraph, 'line_spacing') and paragraph.line_spacing:
                line_height_px = paragraph.line_spacing * 96 / 72

            # space_before (except first paragraph)
            if para_idx > 0:
                space_before = 0
                if hasattr(paragraph, 'space_before') and paragraph.space_before:
                    space_before = paragraph.space_before.pt if hasattr(paragraph.space_before, 'pt') else 0
                else:
                    # PowerPoint default inter-paragraph spacing when not explicit
                    space_before = 2  # ~2pt default gap
                total_height_px += space_before * 96 / 72

            # Paragraph text height
            total_height_px += len(all_wrapped_lines) * line_height_px

            # space_after
            space_after = 0
            if hasattr(paragraph, 'space_after') and paragraph.space_after:
                space_after = paragraph.space_after.pt if hasattr(paragraph.space_after, 'pt') else 0
            total_height_px += space_after * 96 / 72

    # Apply scaled safety factor based on content density.
    # PIL consistently underestimates for dense multi-line content due to
    # PowerPoint's larger line spacing, paragraph gaps, and CJK rendering.
    # Single lines/short text: no correction needed (PowerPoint fits them fine).
    para_count = len([p for p in text_frame.paragraphs if p.text.strip()])
    total_chars = sum(len(p.text) for p in text_frame.paragraphs)

    if para_count <= 2 and total_chars < 50:
        # Short content (titles, labels, badges) — no safety needed
        safety_factor = 1.0
    elif para_count <= 4:
        safety_factor = 1.10
    elif para_count <= 8:
        safety_factor = 1.18
    else:
        # Dense multi-paragraph blocks — highest accumulated error
        safety_factor = 1.25

    total_height_px *= safety_factor

    total_height_inches = total_height_px / 96.0
    overflow_inches = max(0, (total_height_px - usable_height_px) / 96.0)

    return (total_height_inches, overflow_inches)


# ─── Validation Engine ────────────────────────────────────────────────────────

class ValidationIssue:
    """Represents a single validation issue."""

    def __init__(self, slide_num: int, shape_idx: int, issue_type: str,
                 text_preview: str, details: dict):
        self.slide_num = slide_num
        self.shape_idx = shape_idx
        self.issue_type = issue_type
        self.text_preview = text_preview
        self.details = details

    def to_dict(self) -> dict:
        return {
            "slide": self.slide_num,
            "shape_index": self.shape_idx,
            "type": self.issue_type,
            "text_preview": self.text_preview[:80],
            "details": self.details,
        }


def validate_pptx(filepath: str, threshold: float = 0.05) -> List[ValidationIssue]:
    """Validate a PPTX file for text overflow and layout issues.

    Args:
        filepath: Path to the .pptx file
        threshold: Minimum overflow in inches to report (default 0.05")

    Returns:
        List of ValidationIssue objects
    """
    prs = Presentation(filepath)
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)
    issues = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        for shape_idx, shape in enumerate(slide.shapes):
            # Get shape bounds
            shape_left = emu_to_inches(shape.left) if shape.left else 0
            shape_top = emu_to_inches(shape.top) if shape.top else 0
            shape_width = emu_to_inches(shape.width) if shape.width else 0
            shape_height = emu_to_inches(shape.height) if shape.height else 0

            # Check slide boundary overflow
            right_edge = shape_left + shape_width
            bottom_edge = shape_top + shape_height
            if right_edge > slide_width + 0.01:
                issues.append(ValidationIssue(
                    slide_num, shape_idx, "slide_overflow_right",
                    f"Shape extends {right_edge - slide_width:.2f}\" past right edge",
                    {"overflow_inches": round(right_edge - slide_width, 2),
                     "shape_right": round(right_edge, 2),
                     "slide_width": round(slide_width, 2)}
                ))
            if bottom_edge > slide_height + 0.01:
                issues.append(ValidationIssue(
                    slide_num, shape_idx, "slide_overflow_bottom",
                    f"Shape extends {bottom_edge - slide_height:.2f}\" past bottom edge",
                    {"overflow_inches": round(bottom_edge - slide_height, 2),
                     "shape_bottom": round(bottom_edge, 2),
                     "slide_height": round(slide_height, 2)}
                ))

            # Check text overflow
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            full_text = tf.text.strip()
            if not full_text:
                continue

            # Font size check (WARNING only, not a blocking issue)
            # Font size is a SOFT constraint — content visibility takes priority
            # Only warn if font < 7pt (practically unreadable)
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt < 7:
                        issues.append(ValidationIssue(
                            slide_num, shape_idx, "font_too_small",
                            run.text[:60],
                            {"font_size_pt": run.font.size.pt,
                             "minimum_pt": 7,
                             "font_name": run.font.name or "unknown"}
                        ))

            # Measure text overflow
            if shape_width > 0 and shape_height > 0:
                total_h, overflow = measure_text_height(tf, shape_width, shape_height)
                if overflow > threshold:
                    issues.append(ValidationIssue(
                        slide_num, shape_idx, "text_overflow",
                        full_text[:80],
                        {"overflow_inches": round(overflow, 3),
                         "total_text_height": round(total_h, 3),
                         "shape_height": round(shape_height, 3),
                         "shape_width": round(shape_width, 3),
                         "shape_left": round(shape_left, 3),
                         "shape_top": round(shape_top, 3),
                         "suggested_fix": "reduce_font_size" if overflow < 0.3
                                          else "reduce_content_or_expand"}
                    ))

                # Widow line detection: if text wraps and last line < 20% of content
                widow_issues = detect_widow_lines(tf, shape_width, shape_height)
                for widow in widow_issues:
                    issues.append(ValidationIssue(
                        slide_num, shape_idx, "widow_line",
                        widow["text_preview"],
                        widow
                    ))

    return issues


def detect_widow_lines(text_frame, shape_width_inches: float, shape_height_inches: float) -> List[dict]:
    """Detect paragraphs where wrapped text has a widow line (last line < 20% of total chars).

    A "widow line" is when text wraps and the last wrapped line contains only a few
    characters (< 20% of the paragraph's total length), making the layout look bad.
    """
    margin_left = emu_to_inches(text_frame.margin_left) if text_frame.margin_left else 0.1
    margin_right = emu_to_inches(text_frame.margin_right) if text_frame.margin_right else 0.1
    usable_width = shape_width_inches - margin_left - margin_right

    if usable_width <= 0:
        return []

    usable_width_px = int(usable_width * 96)
    dummy_img = Image.new("RGB", (1, 1))
    draw = ImageDraw.Draw(dummy_img)

    widow_issues = []

    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        text = paragraph.text.strip()
        if not text or len(text) < 5:
            continue

        # Get font info
        font_name = "Arial"
        font_size = 14
        if paragraph.runs:
            run = paragraph.runs[0]
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = int(run.font.size.pt)
        elif paragraph.font and paragraph.font.size:
            font_size = int(paragraph.font.size.pt)

        font = load_font(font_name, font_size)

        # Wrap the text
        all_wrapped = []
        for line in text.split("\n"):
            wrapped = wrap_text_line(line, usable_width_px, draw, font)
            all_wrapped.extend(wrapped)

        # Check for widow: multiple lines where last line is < 20% of total chars
        if len(all_wrapped) >= 2:
            total_chars = sum(len(l) for l in all_wrapped)
            last_line_chars = len(all_wrapped[-1].strip())

            if total_chars > 0 and last_line_chars < total_chars * 0.20:
                widow_issues.append({
                    "paragraph_index": para_idx,
                    "total_lines": len(all_wrapped),
                    "last_line_text": all_wrapped[-1].strip(),
                    "last_line_chars": last_line_chars,
                    "total_chars": total_chars,
                    "ratio": round(last_line_chars / total_chars, 3),
                    "font_size_pt": font_size,
                    "shape_width": round(shape_width_inches, 3),
                    "text_preview": text[:80],
                    "suggested_fix": "reduce_font_or_widen",
                })

    return widow_issues


def print_report(issues: List[ValidationIssue], filepath: str) -> None:
    """Print a human-readable report."""
    if not issues:
        print(f"VALIDATION: PASS")
        print(f"FILE: {filepath}")
        print(f"ISSUES: 0")
        return

    print(f"VALIDATION: FAIL")
    print(f"FILE: {filepath}")
    print(f"ISSUES: {len(issues)}")
    print()

    # Group by slide
    by_slide = {}
    for issue in issues:
        by_slide.setdefault(issue.slide_num, []).append(issue)

    for slide_num in sorted(by_slide.keys()):
        slide_issues = by_slide[slide_num]
        print(f"  Slide {slide_num}: {len(slide_issues)} issue(s)")
        for issue in slide_issues:
            if issue.issue_type == "text_overflow":
                d = issue.details
                print(f"    [{issue.issue_type}] overflow={d['overflow_inches']}\""
                      f" (text_h={d['total_text_height']}\", box_h={d['shape_height']}\")")
                print(f"      Text: \"{issue.text_preview}...\"")
                print(f"      Fix: {d.get('suggested_fix', 'manual')}")
            elif issue.issue_type == "font_too_small":
                d = issue.details
                print(f"    [{issue.issue_type}] {d['font_size_pt']}pt < {d['minimum_pt']}pt minimum")
                print(f"      Text: \"{issue.text_preview}\"")
            elif issue.issue_type == "widow_line":
                d = issue.details
                print(f"    [{issue.issue_type}] last line \"{d['last_line_text']}\" "
                      f"= {d['last_line_chars']}/{d['total_chars']} chars ({d['ratio']*100:.0f}%)")
                print(f"      Text: \"{issue.text_preview}...\"")
                print(f"      Fix: {d.get('suggested_fix', 'reduce_font_or_widen')}")
            else:
                print(f"    [{issue.issue_type}] {issue.text_preview}")
        print()


def main():
    parser = argparse.ArgumentParser(description="Validate PPTX for text overflow issues")
    parser.add_argument("input", help="Path to .pptx file to validate")
    parser.add_argument("--output", "-o", help="Path to save JSON report (optional)")
    parser.add_argument("--threshold", type=float, default=0.05,
                        help="Minimum overflow in inches to report (default: 0.05)")
    parser.add_argument("--json", action="store_true",
                        help="Output JSON to stdout instead of human-readable report")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"ERROR: File not found: {args.input}", file=sys.stderr)
        sys.exit(2)

    issues = validate_pptx(args.input, threshold=args.threshold)

    # Output report
    if args.json:
        report = {
            "file": args.input,
            "status": "pass" if not issues else "fail",
            "issue_count": len(issues),
            "issues": [i.to_dict() for i in issues],
        }
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print_report(issues, args.input)

    # Save JSON if requested
    if args.output:
        report = {
            "file": args.input,
            "status": "pass" if not issues else "fail",
            "issue_count": len(issues),
            "issues": [i.to_dict() for i in issues],
        }
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
        print(f"\nReport saved to: {args.output}")

    sys.exit(0 if not issues else 1)


if __name__ == "__main__":
    main()
