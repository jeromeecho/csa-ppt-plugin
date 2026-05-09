#!/usr/bin/env python3
"""Render PPTX slides to PNG images for LLM visual inspection.

Tries multiple rendering backends in order:
1. LibreOffice headless (best quality, full fidelity)
2. python-pptx + PIL fallback (shape-based approximate rendering)

Output: one PNG per slide in the specified output directory.

Usage:
    python render_slides.py input.pptx --output-dir /tmp/slides/ [--slides 55,56,57]
    python render_slides.py input.pptx --output-dir /tmp/slides/ --slides 55-57
"""

import argparse
import os
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("ERROR: python-pptx and Pillow required. pip install python-pptx Pillow", file=sys.stderr)
    sys.exit(2)


def find_libreoffice() -> Optional[str]:
    """Find LibreOffice/soffice binary (cross-platform: macOS/Windows/Linux)."""
    candidates = [
        # Generic (works if on PATH)
        "libreoffice", "soffice",
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        # Linux
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
        "/snap/bin/libreoffice",
        # Windows (common install paths)
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for cmd in candidates:
        try:
            r = subprocess.run([cmd, "--version"], capture_output=True, timeout=5)
            if r.returncode == 0:
                return cmd
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


def render_with_libreoffice(pptx_path: str, output_dir: str, slides: Optional[List[int]] = None) -> List[str]:
    """Render PPTX to PNG via LibreOffice (PDF intermediate)."""
    soffice = find_libreoffice()
    if not soffice:
        return []

    with tempfile.TemporaryDirectory() as tmpdir:
        # Convert PPTX to PDF
        cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path]
        r = subprocess.run(cmd, capture_output=True, timeout=120)
        if r.returncode != 0:
            return []

        pdf_name = Path(pptx_path).stem + ".pdf"
        pdf_path = os.path.join(tmpdir, pdf_name)
        if not os.path.exists(pdf_path):
            return []

        # Convert PDF to PNGs using pdftoppm
        prefix = os.path.join(output_dir, "slide")
        cmd = ["pdftoppm", "-png", "-r", "150", pdf_path, prefix]
        r = subprocess.run(cmd, capture_output=True, timeout=120)
        if r.returncode != 0:
            return []

        # Collect output files
        png_files = sorted(Path(output_dir).glob("slide-*.png"))

        # If specific slides requested, filter
        if slides:
            filtered = []
            for s in slides:
                # pdftoppm names: slide-01.png, slide-02.png, etc.
                for pattern in [f"slide-{s:02d}.png", f"slide-{s:03d}.png", f"slide-{s}.png"]:
                    p = Path(output_dir) / pattern
                    if p.exists():
                        filtered.append(str(p))
                        break
            return filtered

        return [str(p) for p in png_files]


def render_with_pil(pptx_path: str, output_dir: str, slides: Optional[List[int]] = None) -> List[str]:
    """Render PPTX to PNG using python-pptx + PIL (approximate but works without LibreOffice).

    Creates a visual representation of each slide showing:
    - Shape boundaries (rectangles)
    - Text content rendered at approximate positions
    - Background colors
    """
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width / 914400  # inches
    slide_h = prs.slide_height / 914400

    # Render at 150 DPI
    dpi = 150
    img_w = int(slide_w * dpi)
    img_h = int(slide_h * dpi)

    slide_indices = range(len(prs.slides))
    if slides:
        slide_indices = [s - 1 for s in slides if 0 < s <= len(prs.slides)]

    output_files = []

    for slide_idx in slide_indices:
        slide = prs.slides[slide_idx]
        slide_num = slide_idx + 1

        # Create white canvas
        img = Image.new("RGB", (img_w, img_h), "#ffffff")
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            # Get shape position in pixels
            left_px = int((shape.left / 914400) * dpi) if shape.left else 0
            top_px = int((shape.top / 914400) * dpi) if shape.top else 0
            width_px = int((shape.width / 914400) * dpi) if shape.width else 0
            height_px = int((shape.height / 914400) * dpi) if shape.height else 0

            if width_px <= 0 or height_px <= 0:
                continue

            # Draw shape background
            has_fill = False
            fill_color = None
            try:
                fill = shape.fill
                if fill and fill.type is not None:
                    has_fill = True
                    try:
                        fill_color = "#" + str(fill.fore_color.rgb)
                    except:
                        fill_color = "#f0f2f5"
            except:
                pass

            if has_fill and fill_color:
                draw.rectangle(
                    [left_px, top_px, left_px + width_px, top_px + height_px],
                    fill=fill_color, outline="#cccccc"
                )
            else:
                # Just draw outline for shapes with text
                if shape.has_text_frame and shape.text_frame.text.strip():
                    draw.rectangle(
                        [left_px, top_px, left_px + width_px, top_px + height_px],
                        outline="#e0e0e0"
                    )

            # Draw text
            if shape.has_text_frame:
                tf = shape.text_frame
                y_offset = top_px + 4

                for para in tf.paragraphs:
                    text = para.text.strip()
                    if not text:
                        y_offset += 8
                        continue

                    # Get font size
                    font_size = 14
                    font_color = "#333333"
                    is_bold = False
                    for run in para.runs:
                        if run.font.size:
                            font_size = int(run.font.size.pt)
                        if run.font.bold:
                            is_bold = True
                        if run.font.color and run.font.color.rgb:
                            font_color = "#" + str(run.font.color.rgb)
                        break

                    # Scale font for rendering
                    render_size = int(font_size * dpi / 72)

                    # Try to load font
                    font = None
                    for font_path in [
                        "/System/Library/Fonts/PingFang.ttc",
                        "/System/Library/Fonts/Helvetica.ttc",
                        "/Library/Fonts/Arial.ttf",
                    ]:
                        if os.path.exists(font_path):
                            try:
                                font = ImageFont.truetype(font_path, size=render_size)
                                break
                            except:
                                continue

                    if font is None:
                        font = ImageFont.load_default()

                    # Draw text with word wrap
                    max_text_width = width_px - 10
                    lines = text.split("\n") if "\n" in text else [text]

                    for line in lines:
                        if not line:
                            y_offset += render_size + 2
                            continue

                        # Simple word wrap
                        wrapped_lines = []
                        try:
                            line_w = draw.textlength(line, font=font)
                        except:
                            line_w = len(line) * render_size * 0.6

                        if line_w <= max_text_width:
                            wrapped_lines = [line]
                        else:
                            current = ""
                            for char in line:
                                test = current + char
                                try:
                                    tw = draw.textlength(test, font=font)
                                except:
                                    tw = len(test) * render_size * 0.6
                                if tw <= max_text_width:
                                    current = test
                                else:
                                    if current:
                                        wrapped_lines.append(current)
                                    current = char
                            if current:
                                wrapped_lines.append(current)

                        for wline in wrapped_lines:
                            if y_offset + render_size > top_px + height_px:
                                # TEXT OVERFLOW - draw in red to highlight
                                draw.text((left_px + 5, y_offset), wline,
                                         fill="#FF0000", font=font)
                                # Draw red border around shape to indicate overflow
                                draw.rectangle(
                                    [left_px, top_px, left_px + width_px, top_px + height_px],
                                    outline="#FF0000", width=2
                                )
                            else:
                                draw.text((left_px + 5, y_offset), wline,
                                         fill=font_color, font=font)
                            y_offset += render_size + 2

                    y_offset += 4  # paragraph spacing

        # Save
        out_path = os.path.join(output_dir, f"slide-{slide_num:02d}.png")
        img.save(out_path, "PNG")
        output_files.append(out_path)

    return output_files


def parse_slides_arg(slides_str: str, max_slides: int) -> List[int]:
    """Parse slide specification like '55,56,57' or '55-57'."""
    result = []
    for part in slides_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            for s in range(int(start), int(end) + 1):
                if 1 <= s <= max_slides:
                    result.append(s)
        else:
            s = int(part)
            if 1 <= s <= max_slides:
                result.append(s)
    return sorted(set(result))


def main():
    parser = argparse.ArgumentParser(description="Render PPTX slides to PNG for LLM visual inspection")
    parser.add_argument("input", help="Path to .pptx file")
    parser.add_argument("--output-dir", "-o", required=True, help="Directory to save PNG files")
    parser.add_argument("--slides", help="Specific slides to render (e.g., '55,56,57' or '55-57')")
    parser.add_argument("--dpi", type=int, default=150, help="Rendering DPI (default: 150)")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"ERROR: File not found: {args.input}", file=sys.stderr)
        sys.exit(2)

    os.makedirs(args.output_dir, exist_ok=True)

    # Parse slides
    prs = Presentation(args.input)
    slides = None
    if args.slides:
        slides = parse_slides_arg(args.slides, len(prs.slides))
        print(f"Rendering slides: {slides}")
    else:
        print(f"Rendering all {len(prs.slides)} slides")

    # Try LibreOffice first (best quality)
    print("Trying LibreOffice rendering...")
    output_files = render_with_libreoffice(args.input, args.output_dir, slides)

    if output_files:
        print(f"SUCCESS: LibreOffice rendered {len(output_files)} slides")
    else:
        # Fallback to PIL rendering
        print("LibreOffice not available, using PIL fallback rendering...")
        print("(Note: PIL rendering is approximate - text overflow shown in RED)")
        output_files = render_with_pil(args.input, args.output_dir, slides)
        print(f"SUCCESS: PIL rendered {len(output_files)} slides")

    print()
    print("OUTPUT_FILES:")
    for f in output_files:
        print(f"  {f}")
    print()
    print("Next step: Have the LLM read these images to verify:")
    print("  1. All text is fully visible (no clipping/truncation)")
    print("  2. No widow lines (last wrapped line with only 1-2 chars)")
    print("  3. Layout is reasonable (no overlap, no excessive gaps)")


if __name__ == "__main__":
    main()
